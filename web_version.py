import os
import uuid
import threading
import time
from bottle import Bottle, run, request, static_file
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

app = Bottle()

UPLOAD_FOLDER = './uploads'
OUTPUT_FOLDER = './output'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
if not os.path.exists(OUTPUT_FOLDER):
    os.makedirs(OUTPUT_FOLDER)

@app.route('/static/<filepath:path>')
def serve_static(filepath):
    return static_file(filepath, root='static')

@app.route('/')
def upload_form():
    return static_file('upload.html', root='./templates')

@app.route('/upload', method='POST')
def upload():
    user_id = str(uuid.uuid4())  # Unique ID for each user session
    user_upload_folder = os.path.join(UPLOAD_FOLDER, user_id)
    user_output_folder = os.path.join(OUTPUT_FOLDER, user_id)
    os.makedirs(user_upload_folder)
    os.makedirs(user_output_folder)

    files = request.files.getall('upload')
    for file in files:
        if file and file.filename.endswith('.png'):
            file.save(os.path.join(user_upload_folder, file.filename))

    pptx_file = os.path.join(user_output_folder, "output.pptx")
    pngs_to_pptx_16x9(user_upload_folder, pptx_file)

    # Serve the PPTX file for download
    response = static_file("output.pptx", root=user_output_folder, download="output.pptx")

    # Start a thread to delete directories after 5 minutes
    threading.Thread(target=delayed_delete_directory, args=(user_upload_folder, 3000)).start()
    threading.Thread(target=delayed_delete_directory, args=(user_output_folder, 3000)).start()

    return response

def pngs_to_pptx_16x9(png_folder, pptx_file):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for filename in os.listdir(png_folder):
        if filename.endswith(".png"):
            img_path = os.path.join(png_folder, filename)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            img = Image.open(img_path)
            img_width, img_height = img.size

            img_aspect_ratio = img_width / img_height

            if img_aspect_ratio > slide_width / slide_height:
                pic_width = slide_width
                pic_height = slide_width / img_aspect_ratio
            else:
                pic_height = slide_height
                pic_width = slide_height * img_aspect_ratio

            left = (slide_width - pic_width) / 2
            top = (slide_height - pic_height) / 2
            slide.shapes.add_picture(img_path, left, top, width=pic_width, height=pic_height)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = os.path.splitext(filename)[0]

    prs.save(pptx_file)
    print(f"PPTX file saved: {pptx_file}")

@app.route('/download/<user_id>/<filename>')
def download(user_id, filename):
    file_path = os.path.join(OUTPUT_FOLDER, user_id, filename)
    if os.path.isfile(file_path):
        response = static_file(filename, root=os.path.dirname(file_path), download=filename)
        
        # Start a thread to delete the directory after 5 minutes
        threading.Thread(target=delayed_delete_directory, args=(os.path.dirname(file_path), 3000)).start()
        
        return response
    else:
        return "File not found", 404

def delete_directory(directory):
    for root, dirs, files in os.walk(directory, topdown=False):
        for name in files:
            os.remove(os.path.join(root, name))
        for name in dirs:
            os.rmdir(os.path.join(root, name))
    os.rmdir(directory)

def delayed_delete_directory(directory, delay):
    time.sleep(delay)
    if os.path.exists(directory):
        delete_directory(directory)

if __name__ == "__main__":
    run(app, host='localhost', port=8080)
