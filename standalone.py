import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def pngs_to_pptx_16x9(png_folder, pptx_file):
    # Initialize a widescreen presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for filename in os.listdir(png_folder):
        if filename.endswith(".png"):
            img_path = os.path.join(png_folder, filename)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            img = Image.open(img_path)
            img_width, img_height = img.size

            # Calculate aspect ratio and determine the dimensions to fit within slide
            img_aspect_ratio = img_width / img_height

            if img_aspect_ratio > slide_width / slide_height:
                # Fit to slide width
                pic_width = slide_width
                pic_height = slide_width / img_aspect_ratio
            else:
                # Fit to slide height
                pic_height = slide_height
                pic_width = slide_height * img_aspect_ratio

            # Center the image on the slide
            left = (slide_width - pic_width) / 2
            top = (slide_height - pic_height) / 2
            slide.shapes.add_picture(img_path, left, top, width=pic_width, height=pic_height)

            # Optional: Add a title with the filename
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = os.path.splitext(filename)[0]

    prs.save(pptx_file)
    print(f"PPTX file saved: {pptx_file}")

if __name__ == "__main__":
    png_folder = './files'  # Replace with the actual path if needed
    pptx_file = "output.pptx"
    pngs_to_pptx_16x9(png_folder, pptx_file)
